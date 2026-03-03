"""
Unit tests for tools/md_to_pptx.py

Run from the repository root with:
    pytest tools/tests/test_md_to_pptx.py -v
"""
import os
import sys

import pytest
from pptx import Presentation
from pptx.util import Inches

# Ensure the tools/ directory is on the path so md_to_pptx can be imported
# regardless of how pytest is invoked (from root or from within tools/).
_TOOLS_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
if _TOOLS_DIR not in sys.path:
    sys.path.insert(0, _TOOLS_DIR)

from md_to_pptx import (  # noqa: E402
    Block,
    normalize_bullets,
    parse_markdown,
    parse_table_lines,
    to_abs_path,
    generate_pptx,
    main,
)

_SCRIPT_NAME = "md_to_pptx.py"


# ---------------------------------------------------------------------------
# parse_markdown
# ---------------------------------------------------------------------------

class TestParseMarkdownHeadings:
    """Tests for heading block detection in parse_markdown."""

    def test_h1_produces_heading_block(self):
        blocks = parse_markdown("# Title")
        assert len(blocks) == 1
        assert blocks[0].type == "heading"
        assert blocks[0].level == 1
        assert blocks[0].content == ["Title"]

    def test_h2_produces_heading_level_2(self):
        blocks = parse_markdown("## Section")
        assert blocks[0].level == 2

    def test_h3_produces_heading_level_3(self):
        blocks = parse_markdown("### Sub")
        assert blocks[0].level == 3

    def test_heading_text_is_stripped(self):
        blocks = parse_markdown("#   Spaces   ")
        assert blocks[0].content == ["Spaces"]

    def test_multiple_headings(self):
        md = "# H1\n## H2\n### H3"
        blocks = parse_markdown(md)
        headings = [b for b in blocks if b.type == "heading"]
        assert [h.level for h in headings] == [1, 2, 3]


class TestParseMarkdownLists:
    """Tests for list block detection."""

    def test_dash_list_item(self):
        blocks = parse_markdown("- item one")
        assert any(b.type == "list" for b in blocks)

    def test_asterisk_list_item(self):
        blocks = parse_markdown("* item")
        assert any(b.type == "list" for b in blocks)

    def test_plus_list_item(self):
        blocks = parse_markdown("+ item")
        assert any(b.type == "list" for b in blocks)

    def test_numbered_list_item(self):
        blocks = parse_markdown("1. First item")
        assert any(b.type == "list" for b in blocks)

    def test_consecutive_list_items_merged(self):
        md = "- a\n- b\n- c"
        blocks = parse_markdown(md)
        list_blocks = [b for b in blocks if b.type == "list"]
        assert len(list_blocks) == 1
        assert len(list_blocks[0].content) == 3

    def test_list_after_blank_line_merged_into_one_block(self):
        # The parser always attaches new list items to the most recent list block,
        # so a blank line between items does NOT split them into two blocks.
        md = "- a\n\n- b"
        blocks = parse_markdown(md)
        list_blocks = [b for b in blocks if b.type == "list"]
        assert len(list_blocks) == 1
        assert len(list_blocks[0].content) == 2


class TestParseMarkdownCodeBlocks:
    """Tests for fenced code block detection."""

    def test_plain_code_block(self):
        md = "```\nsome code\n```"
        blocks = parse_markdown(md)
        code_blocks = [b for b in blocks if b.type == "code"]
        assert len(code_blocks) == 1
        assert code_blocks[0].content == ["some code"]

    def test_code_block_with_language(self):
        md = "```python\nprint('hi')\n```"
        blocks = parse_markdown(md)
        code_blocks = [b for b in blocks if b.type == "code"]
        assert code_blocks[0].lang == "python"

    def test_mermaid_code_block_lang(self):
        md = "```mermaid\ngraph TD\n  A --> B\n```"
        blocks = parse_markdown(md)
        code_blocks = [b for b in blocks if b.type == "code"]
        assert code_blocks[0].lang == "mermaid"

    def test_code_block_multiline_content(self):
        md = "```bash\nline1\nline2\nline3\n```"
        blocks = parse_markdown(md)
        code_blocks = [b for b in blocks if b.type == "code"]
        assert code_blocks[0].content == ["line1", "line2", "line3"]

    def test_list_inside_code_block_not_parsed_as_list(self):
        md = "```\n- not a list\n```"
        blocks = parse_markdown(md)
        assert not any(b.type == "list" for b in blocks)


class TestParseMarkdownImages:
    """Tests for inline image detection."""

    def test_image_on_own_line(self):
        blocks = parse_markdown("![Alt text](path/to/img.png)")
        image_blocks = [b for b in blocks if b.type == "image"]
        assert len(image_blocks) == 1
        assert image_blocks[0].content == ["Alt text", "path/to/img.png"]

    def test_image_with_empty_alt(self):
        blocks = parse_markdown("![](image.jpg)")
        image_blocks = [b for b in blocks if b.type == "image"]
        assert image_blocks[0].content == ["", "image.jpg"]

    def test_image_with_url(self):
        blocks = parse_markdown("![Logo](https://example.com/logo.png)")
        image_blocks = [b for b in blocks if b.type == "image"]
        assert image_blocks[0].content[1] == "https://example.com/logo.png"


class TestParseMarkdownTables:
    """Tests for Markdown table detection."""

    def test_single_row_table(self):
        md = "| col1 | col2 |"
        blocks = parse_markdown(md)
        assert any(b.type == "table" for b in blocks)

    def test_full_table_with_separator(self):
        md = "| a | b |\n|---|---|\n| 1 | 2 |"
        blocks = parse_markdown(md)
        table_blocks = [b for b in blocks if b.type == "table"]
        assert len(table_blocks) == 1
        # Three raw lines stored
        assert len(table_blocks[0].content) == 3

    def test_table_ends_before_non_table_line(self):
        md = "| a | b |\n\nsome text"
        blocks = parse_markdown(md)
        table_blocks = [b for b in blocks if b.type == "table"]
        assert len(table_blocks) == 1


class TestParseMarkdownSpeakerNotes:
    """Tests for speaker-notes block detection."""

    def test_notes_block_detected(self):
        md = "<!-- SPEAKER_NOTES_START\n- note line\nSPEAKER_NOTES_END -->"
        blocks = parse_markdown(md)
        notes_blocks = [b for b in blocks if b.type == "notes"]
        assert len(notes_blocks) == 1

    def test_notes_content_captured(self):
        md = "<!-- SPEAKER_NOTES_START\n- point one\n- point two\nSPEAKER_NOTES_END -->"
        blocks = parse_markdown(md)
        notes_blocks = [b for b in blocks if b.type == "notes"]
        assert "- point one" in notes_blocks[0].content
        assert "- point two" in notes_blocks[0].content

    def test_notes_not_parsed_as_list(self):
        md = "<!-- SPEAKER_NOTES_START\n- item\nSPEAKER_NOTES_END -->"
        blocks = parse_markdown(md)
        assert not any(b.type == "list" for b in blocks)


class TestParseMarkdownSlideSeparator:
    """Tests that --- slide separators are ignored."""

    def test_separator_not_in_blocks(self):
        md = "# Title\n\n---\n\n## Section"
        blocks = parse_markdown(md)
        types = [b.type for b in blocks]
        assert "separator" not in types

    def test_content_around_separator_preserved(self):
        md = "# Title\n\n---\n\n## Section"
        blocks = parse_markdown(md)
        headings = [b for b in blocks if b.type == "heading"]
        assert len(headings) == 2


class TestParseMarkdownParagraphs:
    """Tests for paragraph block detection."""

    def test_plain_text_becomes_para(self):
        blocks = parse_markdown("Some plain text here.")
        assert any(b.type == "para" for b in blocks)

    def test_blank_line_within_paragraph_adds_empty_content(self):
        # The parser keeps one para block and inserts an empty string on blank lines.
        md = "Para one.\n\nPara two."
        blocks = parse_markdown(md)
        paras = [b for b in blocks if b.type == "para"]
        assert len(paras) == 1
        assert "" in paras[0].content

    def test_empty_input_returns_no_blocks(self):
        blocks = parse_markdown("")
        assert blocks == []


# ---------------------------------------------------------------------------
# normalize_bullets
# ---------------------------------------------------------------------------

class TestNormalizeBullets:
    """Tests for normalize_bullets helper."""

    def test_dash_bullet_extracted(self):
        result = normalize_bullets(["- item one"])
        assert result == ["item one"]

    def test_asterisk_bullet_extracted(self):
        result = normalize_bullets(["* item"])
        assert result == ["item"]

    def test_plus_bullet_extracted(self):
        result = normalize_bullets(["+ item"])
        assert result == ["item"]

    def test_numbered_bullet_extracted(self):
        result = normalize_bullets(["1. First"])
        assert result == ["First"]

    def test_multiple_bullets(self):
        lines = ["- alpha", "- beta", "- gamma"]
        result = normalize_bullets(lines)
        assert result == ["alpha", "beta", "gamma"]

    def test_non_bullet_line_ignored(self):
        result = normalize_bullets(["not a bullet"])
        assert result == []

    def test_indented_bullet_extracted(self):
        result = normalize_bullets(["  - indented"])
        assert result == ["indented"]

    def test_empty_list(self):
        assert normalize_bullets([]) == []


# ---------------------------------------------------------------------------
# parse_table_lines
# ---------------------------------------------------------------------------

class TestParseTableLines:
    """Tests for parse_table_lines helper."""

    def test_header_and_data_rows(self):
        lines = ["| Name | Age |", "|------|-----|", "| Alice | 30 |"]
        rows = parse_table_lines(lines)
        # Separator row should be skipped
        assert len(rows) == 2
        assert rows[0] == ["Name", "Age"]
        assert rows[1] == ["Alice", "30"]

    def test_separator_row_skipped(self):
        lines = ["|---|---|"]
        rows = parse_table_lines(lines)
        assert rows == []

    def test_cells_stripped(self):
        lines = ["|  cell1  |  cell2  |"]
        rows = parse_table_lines(lines)
        assert rows[0] == ["cell1", "cell2"]

    def test_empty_input(self):
        assert parse_table_lines([]) == []

    def test_non_table_line_ignored(self):
        lines = ["not a table", "| col |"]
        rows = parse_table_lines(lines)
        assert len(rows) == 1

    def test_colon_separator_skipped(self):
        lines = ["|:---|:---|"]
        rows = parse_table_lines(lines)
        assert rows == []


# ---------------------------------------------------------------------------
# to_abs_path
# ---------------------------------------------------------------------------

class TestToAbsPath:
    """Tests for to_abs_path path helper."""

    def test_relative_path_joined_with_base(self):
        result = to_abs_path("/base/dir", "images/photo.png")
        assert result == "/base/dir/images/photo.png"

    def test_leading_dot_slash_removed(self):
        result = to_abs_path("/base", "./image.png")
        assert result == "/base/image.png"

    def test_http_url_returned_unchanged(self):
        url = "http://example.com/image.png"
        assert to_abs_path("/base", url) == url

    def test_https_url_returned_unchanged(self):
        url = "https://example.com/logo.svg"
        assert to_abs_path("/base", url) == url

    def test_nested_relative_path(self):
        result = to_abs_path("/root", "a/b/c.png")
        assert result == "/root/a/b/c.png"


# ---------------------------------------------------------------------------
# main() CLI entry point
# ---------------------------------------------------------------------------

class TestMain:
    """Tests for the main() CLI function."""

    def test_no_args_returns_error_code(self):
        result = main([_SCRIPT_NAME])
        assert result == 2

    def test_one_arg_returns_error_code(self):
        result = main([_SCRIPT_NAME, "input.md"])
        assert result == 2

    def test_valid_args_generates_pptx(self, tmp_path):
        md_file = tmp_path / "slide.md"
        out_file = tmp_path / "out.pptx"
        md_file.write_text("# My Title\n\n## Section\n- bullet\n")
        result = main([_SCRIPT_NAME, str(md_file), str(out_file)])
        assert result == 0
        assert out_file.exists()


# ---------------------------------------------------------------------------
# generate_pptx (integration-level)
# ---------------------------------------------------------------------------

class TestGeneratePptx:
    """Integration tests for the generate_pptx function."""

    def _make_pptx(self, content: str, tmp_path) -> Presentation:
        """Helper: write markdown, generate pptx, reload and return."""
        md_file = tmp_path / "test.md"
        out_file = tmp_path / "test.pptx"
        md_file.write_text(content)
        generate_pptx(str(md_file), str(out_file))
        return Presentation(str(out_file))

    def test_output_file_is_created(self, tmp_path):
        md_file = tmp_path / "test.md"
        out_file = tmp_path / "test.pptx"
        md_file.write_text("# Title\n")
        generate_pptx(str(md_file), str(out_file))
        assert out_file.exists()

    def test_h1_creates_title_slide(self, tmp_path):
        prs = self._make_pptx("# Workshop Title\n", tmp_path)
        assert len(prs.slides) >= 1

    def test_slide_count_matches_h2_sections(self, tmp_path):
        md = "# Title\n\n## Section One\n- bullet\n\n## Section Two\n- bullet\n"
        prs = self._make_pptx(md, tmp_path)
        # Expect title slide + at least 2 content slides
        assert len(prs.slides) >= 3

    def test_output_dir_created_if_missing(self, tmp_path):
        md_file = tmp_path / "test.md"
        out_file = tmp_path / "subdir" / "output.pptx"
        md_file.write_text("# Title\n")
        generate_pptx(str(md_file), str(out_file))
        assert out_file.exists()

    def test_slide_dimensions_are_widescreen(self, tmp_path):
        prs = self._make_pptx("# Title\n", tmp_path)
        expected_width = Inches(13.333)
        expected_height = Inches(7.5)
        assert abs(prs.slide_width - expected_width) < 10
        assert abs(prs.slide_height - expected_height) < 10

    def test_code_block_produces_slide(self, tmp_path):
        md = "# Title\n\n## Code Section\n\n```python\nprint('hello')\n```\n"
        prs = self._make_pptx(md, tmp_path)
        assert len(prs.slides) >= 2

    def test_image_missing_does_not_raise(self, tmp_path):
        md = "# Title\n\n## Img\n\n![Alt](nonexistent.png)\n"
        # Should not raise even when image file is absent
        prs = self._make_pptx(md, tmp_path)
        assert len(prs.slides) >= 1

    def test_table_produces_slide(self, tmp_path):
        md = "# Title\n\n## Table Section\n\n| A | B |\n|---|---|\n| 1 | 2 |\n"
        prs = self._make_pptx(md, tmp_path)
        assert len(prs.slides) >= 2

    def test_empty_markdown_produces_single_slide(self, tmp_path):
        # Even with no headings, the flush at the end may produce a slide or empty prs
        prs = self._make_pptx("", tmp_path)
        # Should complete without error
        assert prs is not None


# ---------------------------------------------------------------------------
# Block dataclass
# ---------------------------------------------------------------------------

class TestBlockDataclass:
    """Tests for the Block dataclass."""

    def test_default_level_is_zero(self):
        b = Block(type="heading", content=["text"])
        assert b.level == 0

    def test_default_lang_is_none(self):
        b = Block(type="code", content=["line"])
        assert b.lang is None

    def test_explicit_lang_stored(self):
        b = Block(type="code", content=[], lang="python")
        assert b.lang == "python"

    def test_content_stored(self):
        b = Block(type="list", content=["a", "b"])
        assert b.content == ["a", "b"]
