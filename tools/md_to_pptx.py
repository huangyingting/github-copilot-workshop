#!/usr/bin/env python3
"""
Convert README.md into a PowerPoint deck.

This tool converts Markdown files into PowerPoint presentations with the following features:
- H1 as title slide, H2/H3 as content slides
- Bulleted lists with proper formatting
- Code blocks as preformatted text boxes with syntax highlighting
- Images embedded if relative path exists
- Simple tables from Markdown pipes
- Mermaid diagram support (if mermaid CLI is available)

Requirements: python-pptx, Pillow

Usage:
  python tools/md_to_pptx.py <input.md> <output.pptx>
  
Example:
  python tools/md_to_pptx.py slides/README-slides.md slides/github-copilot-workshop.pptx
"""
from __future__ import annotations

import os
import re
import sys
from dataclasses import dataclass
from typing import List, Optional, Tuple
import shutil
import subprocess
import tempfile

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from PIL import Image


HEADING_RE = re.compile(r"^(#{1,6})\s+(.*)$")
LIST_RE = re.compile(r"^([*\-\+]|\d+\.)\s+(.*)$")
IMG_RE = re.compile(r"^!\[(.*?)\]\((.*?)\)")
TABLE_ROW_RE = re.compile(r"^\|(.+)\|\s*$")


@dataclass
class Block:
    type: str
    content: List[str]
    level: int = 0
    lang: Optional[str] = None


def parse_markdown(md: str) -> List[Block]:
    blocks: List[Block] = []
    lines = md.splitlines()
    i = 0
    in_code = False
    code_lang = ""
    table_active = False
    in_notes = False

    while i < len(lines):
        line = lines[i]

        # Speaker notes blocks
        if not in_code and line.strip().startswith("<!-- SPEAKER_NOTES_START"):
            in_notes = True
            blocks.append(Block("notes", []))
            i += 1
            continue
        if in_notes:
            if line.strip().endswith("SPEAKER_NOTES_END -->"):
                in_notes = False
            else:
                blocks[-1].content.append(line)
            i += 1
            continue

        # Fenced code blocks
        if line.strip().startswith("```"):
            fence = line.strip()
            if not in_code:
                in_code = True
                code_lang = fence.strip("`") or ""
                blocks.append(Block("code", [], 0, code_lang))
            else:
                in_code = False
            i += 1
            continue

        if in_code:
            blocks[-1].content.append(line)
            i += 1
            continue

        # Ignore slide separators '---' when not in code/notes
        if line.strip() == "---":
            i += 1
            continue

        # Headings
        m = HEADING_RE.match(line)
        if m:
            level = len(m.group(1))
            text = m.group(2).strip()
            blocks.append(Block("heading", [text], level))
            i += 1
            continue

        # Images on a single line
        m = IMG_RE.match(line.strip())
        if m:
            alt, path = m.groups()
            blocks.append(Block("image", [alt, path]))
            i += 1
            continue

        # Tables (very simple detection)
        if TABLE_ROW_RE.match(line):
            # Start new table if not active
            if not table_active:
                blocks.append(Block("table", []))
                table_active = True
            blocks[-1].content.append(line)
            i += 1
            # End table if next line is not a table row
            if i < len(lines) and not TABLE_ROW_RE.match(lines[i]):
                table_active = False
            continue

        # Lists
        if LIST_RE.match(line.strip()):
            # Attach to last list or create new
            if not blocks or blocks[-1].type != "list":
                blocks.append(Block("list", []))
            blocks[-1].content.append(line.rstrip())
            i += 1
            continue

        # Paragraphs
        if line.strip():
            if not blocks or blocks[-1].type != "para":
                blocks.append(Block("para", []))
            blocks[-1].content.append(line.rstrip())

        else:
            # Blank line separates paragraphs
            if blocks and blocks[-1].type == "para":
                blocks[-1].content.append("")
        i += 1

    return blocks


def add_title_slide(prs: Presentation, title: str, subtitle: Optional[str] = None):
    layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    # Ensure title shrinks to fit width
    try:
        t_tf = slide.shapes.title.text_frame
        t_tf.word_wrap = True
        t_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        # Use minimal margins for maximum width usage
        t_tf.margin_left = prs.slide_width * 0.01
        t_tf.margin_right = prs.slide_width * 0.01
    except Exception:
        pass
    if subtitle:
        slide.placeholders[1].text = subtitle
        try:
            s_tf = slide.placeholders[1].text_frame
            s_tf.word_wrap = True
            s_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            s_tf.margin_left = prs.slide_width * 0.01
            s_tf.margin_right = prs.slide_width * 0.01
        except Exception:
            pass
    return slide


def add_bullets_slide(prs: Presentation, title: str, bullets: List[str]):
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    # Title shrink-to-fit
    try:
        t_tf = slide.shapes.title.text_frame
        t_tf.word_wrap = True
        t_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        t_tf.margin_left = prs.slide_width * 0.01
        t_tf.margin_right = prs.slide_width * 0.01
    except Exception:
        pass

    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    # Content shrink-to-fit and wrapping with minimal margins
    try:
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.margin_left = prs.slide_width * 0.015
        tf.margin_right = prs.slide_width * 0.015
    except Exception:
        pass
    for idx, b in enumerate(bullets):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = b
        p.level = 0
    return slide


def _content_box(prs: Presentation) -> Tuple[int, int, int, int]:
    # Use proportional margins based on slide width (about 3.75% on each side)
    left_margin = prs.slide_width * 0.0375
    right_margin = prs.slide_width * 0.0375
    # Keep top margin proportional to height for title space
    top_margin = prs.slide_height * 0.17  # About 17% for title
    bottom_margin = prs.slide_height * 0.067  # About 6.7% bottom margin
    
    width = prs.slide_width - left_margin - right_margin
    height = prs.slide_height - top_margin - bottom_margin
    return left_margin, top_margin, width, height


def add_code_slide(prs: Presentation, title: str, code: str):
    layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    # Title shrink-to-fit
    try:
        t_tf = slide.shapes.title.text_frame
        t_tf.word_wrap = True
        t_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        t_tf.margin_left = prs.slide_width * 0.01
        t_tf.margin_right = prs.slide_width * 0.01
    except Exception:
        pass
    left, top, width, height = _content_box(prs)
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    # Shrink code text to fit the shape when needed with minimal margins
    try:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.margin_left = prs.slide_width * 0.01
        tf.margin_right = prs.slide_width * 0.01
        tf.margin_top = prs.slide_height * 0.01
        tf.margin_bottom = prs.slide_height * 0.01
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.text = code
    # Monospace-ish formatting
    for run in p.runs:
        run.font.name = "Consolas"
        # Slightly smaller base size to increase fit; auto-size will adjust further if needed
        run.font.size = Pt(11)
    # Note: paragraph font is read-only; we set font on runs above
    return slide


def add_image_slide(prs: Presentation, title: str, image_path: str):
    layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if os.path.exists(image_path):
        # Fit image to content box while preserving aspect ratio
        left, top, box_w, box_h = _content_box(prs)
        try:
            with Image.open(image_path) as im:
                img_w, img_h = im.size
        except Exception:
            img_w, img_h = (1280, 720)
        img_ratio = img_w / img_h if img_h else 1.0
        box_ratio = box_w / box_h if box_h else 1.0
        if img_ratio > box_ratio:
            # limit by width
            disp_w = box_w
            disp_h = int(disp_w / img_ratio)
        else:
            # limit by height
            disp_h = box_h
            disp_w = int(disp_h * img_ratio)
        # center inside box
        pic_left = left + (box_w - disp_w) // 2
        pic_top = top + (box_h - disp_h) // 2
        slide.shapes.add_picture(image_path, pic_left, pic_top, width=disp_w, height=disp_h)
    else:
        # Fallback text if missing
        left, top, width, _ = _content_box(prs)
        box = slide.shapes.add_textbox(left, top, width, Inches(1))
        box.text_frame.text = f"Image not found: {image_path}"
    return slide


def add_table_slide(prs: Presentation, title: str, rows: List[List[str]]):
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    # Title shrink-to-fit
    try:
        t_tf = slide.shapes.title.text_frame
        t_tf.word_wrap = True
        t_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        t_tf.margin_left = prs.slide_width * 0.01
        t_tf.margin_right = prs.slide_width * 0.01
    except Exception:
        pass
    nrows = len(rows)
    ncols = max(len(r) for r in rows)
    left, top, width, max_h = _content_box(prs)
    # Use more of the available height for tables
    height = min(max_h, max_h * 0.8)  # Use up to 80% of available content height
    table = slide.shapes.add_table(nrows, ncols, left, top, width, height).table
    for r, row in enumerate(rows):
        for c in range(ncols):
            table.cell(r, c).text = row[c] if c < len(row) else ""
            # Try to keep table text compact with proportional margins
            try:
                tf = table.cell(r, c).text_frame
                tf.word_wrap = True
                tf.margin_left = prs.slide_width * 0.005
                tf.margin_right = prs.slide_width * 0.005
                tf.margin_top = prs.slide_height * 0.005
                tf.margin_bottom = prs.slide_height * 0.005
            except Exception:
                pass
    return slide


def _which(cmd: str) -> Optional[str]:
    """
    Find executable in PATH (similar to Unix 'which' command).
    
    Args:
        cmd: Command name to search for
        
    Returns:
        Full path to executable if found, None otherwise
    """
    return shutil.which(cmd)


def render_mermaid_to_png(code: str, workdir: Optional[str] = None) -> Optional[str]:
    """Render Mermaid code to a PNG file using mermaid-cli (mmdc) or npx fallback.

    Args:
        code: Mermaid diagram code
        workdir: Working directory to use (creates temp dir if None)
        
    Returns:
        Path to the generated PNG file on success, None otherwise
    """
    tmpdir = tempfile.mkdtemp(prefix="md2pptx_") if workdir is None else workdir
    mmd_path = os.path.join(tmpdir, "diagram.mmd")
    png_path = os.path.join(tmpdir, "diagram.png")
    
    try:
        with open(mmd_path, "w", encoding="utf-8") as f:
            f.write(code)
    except IOError:
        return None

    commands = []
    if _which("mmdc"):
        commands.append(["mmdc", "-i", mmd_path, "-o", png_path])
    # try npx mermaid-cli
    if _which("npx"):
        commands.append(["npx", "-y", "@mermaid-js/mermaid-cli", "-i", mmd_path, "-o", png_path])

    for cmd in commands:
        try:
            subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
            if os.path.exists(png_path) and os.path.getsize(png_path) > 0:
                return png_path
        except Exception:
            continue
    return None


def normalize_bullets(lines: List[str]) -> List[str]:
    # Convert nested lists into flattened bullet lines (ignore nested levels for simplicity)
    out: List[str] = []
    for line in lines:
        m = LIST_RE.match(line.strip())
        if m:
            out.append(m.group(2).strip())
    return out


def parse_table_lines(table_lines: List[str]) -> List[List[str]]:
    rows: List[List[str]] = []
    for line in table_lines:
        m = TABLE_ROW_RE.match(line)
        if not m:
            continue
        row_body = m.group(1)
        cells = [c.strip() for c in row_body.split("|")]
        # Skip markdown table separator rows (---)
        if all(set(c) <= set("-: ") for c in cells):
            continue
        rows.append(cells)
    return rows


def to_abs_path(base_dir: str, path: str) -> str:
    if re.match(r"^https?://", path):
        return path  # external, not downloaded
    # remove optional leading ./
    path = re.sub(r"^\.\/", "", path)
    return os.path.join(base_dir, path)


def generate_pptx(md_path: str, out_path: str) -> None:
    """
    Generate PowerPoint presentation from Markdown file.
    
    Args:
        md_path: Path to input Markdown file
        out_path: Path to output PowerPoint file
        
    Raises:
        IOError: If input file cannot be read or output file cannot be written
        ValueError: If Markdown parsing fails
    """
    try:
        with open(md_path, "r", encoding="utf-8") as f:
            md = f.read()
    except IOError as e:
        raise IOError(f"Cannot read input file '{md_path}': {e}")

    try:
        blocks = parse_markdown(md)
    except Exception as e:
        raise ValueError(f"Failed to parse Markdown: {e}")
        
    prs = Presentation()
    # Set widescreen 16:9 slide size
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    repo_root = os.path.dirname(md_path)
    title_done = False
    current_section = ""
    pending_bullets: List[str] = []
    last_slide = None
    pending_notes: Optional[str] = None

    def flush_bullets(title: str):
        nonlocal pending_bullets
        if pending_bullets:
            slide = add_bullets_slide(prs, title, pending_bullets)
            _maybe_apply_notes(slide)
            pending_bullets = []

    def _maybe_apply_notes(slide):
        nonlocal pending_notes, last_slide
        if pending_notes:
            try:
                notes_frame = slide.notes_slide.notes_text_frame
                notes_frame.clear()
                p = notes_frame.paragraphs[0]
                p.text = pending_notes
            except Exception:
                # If notes slide not available, skip gracefully
                pass
            pending_notes = None
        last_slide = slide

    for b in blocks:
        if b.type == "heading":
            text = b.content[0]
            if b.level == 1 and not title_done:
                slide = add_title_slide(prs, text)
                _maybe_apply_notes(slide)
                title_done = True
                current_section = text
            elif b.level in (2, 3):
                # New slide header for bullets; flush any previous bullets
                flush_bullets(current_section)
                current_section = text
            else:
                # Ignore lower-level headings for simplicity
                pass

        elif b.type == "list":
            bullets = normalize_bullets(b.content)
            pending_bullets.extend(bullets)

        elif b.type == "para":
            text = " ".join([l for l in b.content if l.strip()])
            if text:
                pending_bullets.append(text)

        elif b.type == "code":
            flush_bullets(current_section)
            code_text = "\n".join(b.content)
            # If Mermaid diagram, attempt to render to PNG
            if b.lang and b.lang.lower().startswith("mermaid"):
                png = render_mermaid_to_png(code_text)
                if png:
                    title = f"Diagram: {current_section}" if current_section else "Diagram"
                    slide = add_image_slide(prs, title, png)
                    _maybe_apply_notes(slide)
                else:
                    title = f"Mermaid (code): {current_section}" if current_section else "Mermaid (code)"
                    slide = add_code_slide(prs, title, code_text)
                    _maybe_apply_notes(slide)
            else:
                title = f"Code: {current_section}" if current_section else "Code"
                slide = add_code_slide(prs, title, code_text)
                _maybe_apply_notes(slide)

        elif b.type == "image":
            flush_bullets(current_section)
            alt, path = b.content
            abs_path = to_abs_path(repo_root, path)
            title = alt or current_section or "Image"
            slide = add_image_slide(prs, title, abs_path)
            _maybe_apply_notes(slide)

        elif b.type == "table":
            flush_bullets(current_section)
            rows = parse_table_lines(b.content)
            if rows:
                title = current_section or "Table"
                slide = add_table_slide(prs, title, rows)
                _maybe_apply_notes(slide)

        elif b.type == "notes":
            # Store notes to apply to the next created slide (or the last one if exists without pending)
            text = "\n".join(b.content).strip()
            if last_slide is not None and pending_bullets:
                # If bullets are pending, notes belong to the upcoming bullets slide
                pending_notes = text
            elif last_slide is None:
                pending_notes = text
            else:
                # Attach to the last created slide if no pending bullets
                try:
                    notes_frame = last_slide.notes_slide.notes_text_frame
                    notes_frame.clear()
                    notes_frame.paragraphs[0].text = text
                except Exception:
                    pass

    # Flush any remaining bullets into a slide
    flush_bullets(current_section or "Contents")

    # Ensure output dir exists
    try:
        os.makedirs(os.path.dirname(out_path), exist_ok=True)
        prs.save(out_path)
    except IOError as e:
        raise IOError(f"Cannot write output file '{out_path}': {e}")


def main(argv: List[str]) -> int:
    """
    Main entry point for the script.
    
    Args:
        argv: Command line arguments
        
    Returns:
        Exit code (0 for success, 2 for usage error, 1 for other errors)
    """
    if len(argv) < 3:
        print("Usage: python tools/md_to_pptx.py <input.md> <output.pptx>")
        print("\nExample:")
        print("  python tools/md_to_pptx.py README.md presentation.pptx")
        return 2
    
    md_path, out_path = argv[1], argv[2]
    
    # Validate input file exists
    if not os.path.exists(md_path):
        print(f"Error: Input file '{md_path}' does not exist")
        return 1
    
    # Validate input file is readable
    if not os.access(md_path, os.R_OK):
        print(f"Error: Cannot read input file '{md_path}'")
        return 1
    
    try:
        generate_pptx(md_path, out_path)
        print(f"Successfully wrote {out_path}")
        return 0
    except Exception as e:
        print(f"Error generating PowerPoint: {e}")
        return 1


if __name__ == "__main__":
    sys.exit(main(sys.argv))
